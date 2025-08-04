#!/usr/bin/env python3
"""
News to PowerPoint Converter

This script takes news data from the news aggregator and creates a professional
PowerPoint presentation with the latest news.
"""

import json
import os
from datetime import datetime
from typing import List, Dict
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class NewsToPPTConverter:
    """Convert news data to PowerPoint presentation"""
    
    def __init__(self):
        self.presentation = Presentation()
        self.slide_width = Inches(10)
        self.slide_height = Inches(7.5)
        
        # Color scheme
        self.colors = {
            'primary': RGBColor(31, 73, 125),      # Dark blue
            'secondary': RGBColor(79, 129, 189),   # Light blue
            'accent': RGBColor(192, 80, 77),       # Red
            'text': RGBColor(64, 64, 64),          # Dark gray
            'light_gray': RGBColor(242, 242, 242), # Light gray
            'white': RGBColor(255, 255, 255)       # White
        }

    def create_title_slide(self, news_data: Dict):
        """Create the title slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "📰 Latest News Summary"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.color.rgb = self.colors['primary']
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle = slide.placeholders[1]
        date_str = news_data.get('date', datetime.now().strftime('%Y-%m-%d'))
        formatted_date = datetime.strptime(date_str, '%Y-%m-%d').strftime('%B %d, %Y')
        
        subtitle.text = f"Daily News Digest - {formatted_date}\n"
        subtitle.text += f"📊 {news_data.get('total_articles', 0)} Articles Aggregated\n"
        subtitle.text += f"🕒 Generated at {datetime.now().strftime('%I:%M %p')}"
        
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(18)
        subtitle_paragraph.font.color.rgb = self.colors['secondary']
        subtitle_paragraph.alignment = PP_ALIGN.CENTER

    def create_summary_slide(self, news_data: Dict):
        """Create a summary slide with statistics"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "📊 News Summary Overview"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.color.rgb = self.colors['primary']
        
        # Content
        content = slide.placeholders[1]
        articles = news_data.get('articles', [])
        
        # Count articles by source
        source_counts = {}
        sports_count = 0
        tech_count = 0
        world_news_count = 0
        
        for article in articles:
            source = article.get('source', 'Unknown')
            source_counts[source] = source_counts.get(source, 0) + 1
            
            # Categorize articles
            if 'Sports' in source or 'Cricket' in source:
                sports_count += 1
            elif 'Hacker News' in source:
                tech_count += 1
            elif 'WorldNews' in source:
                world_news_count += 1
        
        summary_text = f"📈 Total Articles: {len(articles)}\n\n"
        summary_text += "📰 Sources Breakdown:\n"
        for source, count in sorted(source_counts.items(), key=lambda x: x[1], reverse=True):
            summary_text += f"   • {source}: {count} articles\n"
        
        summary_text += f"\n🏷️ Category Breakdown:\n"
        summary_text += f"   🏏 Sports: {sports_count} articles\n"
        summary_text += f"   💻 Technology: {tech_count} articles\n"
        summary_text += f"   🌍 World News: {world_news_count} articles\n"
        summary_text += f"   📰 Other: {len(articles) - sports_count - tech_count - world_news_count} articles"
        
        content.text = summary_text
        content_paragraph = content.text_frame.paragraphs[0]
        content_paragraph.font.size = Pt(16)
        content_paragraph.font.color.rgb = self.colors['text']

    def create_category_slide(self, title: str, articles: List[Dict], emoji: str = "📰"):
        """Create a slide for a specific category of news"""
        slide_layout = self.presentation.slide_layouts[1]  # Title and content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        slide_title = slide.shapes.title
        slide_title.text = f"{emoji} {title}"
        title_paragraph = slide_title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.color.rgb = self.colors['primary']
        
        # Content
        content = slide.placeholders[1]
        content_text = ""
        
        for i, article in enumerate(articles[:8], 1):  # Limit to 8 articles per slide
            title_text = article.get('title', 'No Title')
            if len(title_text) > 80:
                title_text = title_text[:77] + "..."
            
            source = article.get('source', 'Unknown')
            
            content_text += f"{i}. {title_text}\n"
            content_text += f"   📍 Source: {source}\n"
            
            # Add publication time if available
            pub_time = article.get('published_at', '')
            if pub_time:
                try:
                    if 'T' in pub_time:
                        dt = datetime.fromisoformat(pub_time.replace('Z', '+00:00'))
                        time_str = dt.strftime('%I:%M %p')
                        content_text += f"   🕒 {time_str}\n"
                except:
                    pass
            
            content_text += "\n"
        
        if len(articles) > 8:
            content_text += f"... and {len(articles) - 8} more articles"
        
        content.text = content_text
        content_paragraph = content.text_frame.paragraphs[0]
        content_paragraph.font.size = Pt(14)
        content_paragraph.font.color.rgb = self.colors['text']

    def create_top_stories_slide(self, articles: List[Dict]):
        """Create a slide with top stories"""
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "🔥 Top Stories of the Day"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.color.rgb = self.colors['primary']
        
        # Content - Top 5 stories with more details
        content = slide.placeholders[1]
        content_text = ""
        
        for i, article in enumerate(articles[:5], 1):
            title_text = article.get('title', 'No Title')
            if len(title_text) > 70:
                title_text = title_text[:67] + "..."
            
            description = article.get('description', '')
            if description and len(description) > 100:
                description = description[:97] + "..."
            
            source = article.get('source', 'Unknown')
            
            content_text += f"{i}. {title_text}\n"
            if description and not description.startswith('Score:'):
                content_text += f"   📝 {description}\n"
            content_text += f"   📍 {source}\n\n"
        
        content.text = content_text
        content_paragraph = content.text_frame.paragraphs[0]
        content_paragraph.font.size = Pt(14)
        content_paragraph.font.color.rgb = self.colors['text']

    def create_sports_slide(self, articles: List[Dict]):
        """Create a dedicated slide for Indian sports news"""
        sports_articles = [
            article for article in articles 
            if 'Sports' in article.get('source', '') or 'Cricket' in article.get('source', '') or 
               article.get('category') in ['Indian Sports', 'Cricket']
        ]
        
        if sports_articles:
            self.create_category_slide("Indian Sports News", sports_articles, "🏏")

    def create_tech_slide(self, articles: List[Dict]):
        """Create a dedicated slide for technology news"""
        tech_articles = [
            article for article in articles 
            if 'Hacker News' in article.get('source', '')
        ]
        
        if tech_articles:
            self.create_category_slide("Technology & Startup News", tech_articles, "💻")

    def create_world_news_slide(self, articles: List[Dict]):
        """Create a dedicated slide for world news"""
        world_articles = [
            article for article in articles 
            if 'WorldNews' in article.get('source', '')
        ]
        
        if world_articles:
            self.create_category_slide("World News", world_articles, "🌍")

    def create_closing_slide(self):
        """Create a closing slide"""
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Thank You!"
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.color.rgb = self.colors['primary']
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = "📰 Stay Informed, Stay Updated\n\n"
        subtitle.text += "Generated by News Aggregator\n"
        subtitle.text += f"🕒 {datetime.now().strftime('%B %d, %Y at %I:%M %p')}"
        
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(18)
        subtitle_paragraph.font.color.rgb = self.colors['secondary']
        subtitle_paragraph.alignment = PP_ALIGN.CENTER

    def convert_news_to_ppt(self, json_file: str, output_file: str = None):
        """Convert news JSON to PowerPoint presentation"""
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                news_data = json.load(f)
        except FileNotFoundError:
            print(f"❌ Error: File '{json_file}' not found.")
            return False
        except json.JSONDecodeError:
            print(f"❌ Error: Invalid JSON in file '{json_file}'.")
            return False
        
        articles = news_data.get('articles', [])
        if not articles:
            print("❌ No articles found in the news data.")
            return False
        
        print(f"🔄 Converting {len(articles)} articles to PowerPoint...")
        
        # Create slides
        self.create_title_slide(news_data)
        self.create_summary_slide(news_data)
        self.create_top_stories_slide(articles)
        self.create_sports_slide(articles)
        self.create_tech_slide(articles)
        self.create_world_news_slide(articles)
        self.create_closing_slide()
        
        # Save presentation
        if not output_file:
            date_str = news_data.get('date', datetime.now().strftime('%Y-%m-%d'))
            output_file = f"news_presentation_{date_str.replace('-', '_')}.pptx"
        
        try:
            self.presentation.save(output_file)
            print(f"✅ PowerPoint presentation saved as '{output_file}'")
            print(f"📊 Created {len(self.presentation.slides)} slides")
            return True
        except Exception as e:
            print(f"❌ Error saving presentation: {e}")
            return False


def main():
    """Main function"""
    parser = argparse.ArgumentParser(description='Convert news JSON to PowerPoint presentation')
    parser.add_argument('json_file', nargs='?', 
                       help='JSON file containing news data (default: latest news_*.json)')
    parser.add_argument('--output', '-o', type=str,
                       help='Output PowerPoint filename')
    
    args = parser.parse_args()
    
    # Find JSON file if not specified
    json_file = args.json_file
    if not json_file:
        # Look for the most recent news JSON file
        json_files = [f for f in os.listdir('.') if f.startswith('news_') and f.endswith('.json')]
        if json_files:
            json_file = sorted(json_files)[-1]  # Get the most recent
            print(f"📁 Using latest news file: {json_file}")
        else:
            print("❌ No news JSON file found. Please run the news aggregator first:")
            print("   python news_aggregator.py --save")
            return
    
    # Convert to PowerPoint
    converter = NewsToPPTConverter()
    success = converter.convert_news_to_ppt(json_file, args.output)
    
    if success:
        print("\n🎉 Conversion completed successfully!")
        print("💡 You can now open the PowerPoint file to view your news presentation.")
    else:
        print("\n❌ Conversion failed. Please check the error messages above.")


if __name__ == "__main__":
    main()
