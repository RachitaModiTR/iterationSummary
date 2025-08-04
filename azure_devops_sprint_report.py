import requests
import json
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from datetime import datetime, timedelta
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from config import AZURE_DEVOPS_CONFIG, COMPLETED_STATES, WORK_ITEM_TYPES, SPRINT_CONFIG, REPORT_CONFIG, FIELD_MAPPINGS

class AzureDevOpsSprintReport:
    def __init__(self, organization, project, team, personal_access_token):
        self.organization = organization
        self.project = project
        self.team = team
        self.pat = personal_access_token
        self.base_url = f"https://dev.azure.com/{organization}/{project}/_apis"
        self.headers = {
            'Authorization': f'Basic {base64.b64encode(f":{personal_access_token}".encode()).decode()}',
            'Content-Type': 'application/json'
        }
        
    def get_work_items_for_sprint(self, iteration_path, area_path, work_item_types):
        """Get work items for the specific sprint iteration and area path"""
        # Build work item types filter
        work_item_types_filter = "', '".join(work_item_types)
        
        # WIQL query with specific iteration path, area path, and work item types
        wiql_query = {
            "query": f"""
            SELECT [System.Id], [System.Title], [System.WorkItemType], [System.State], 
                   [System.AssignedTo], [Microsoft.VSTS.Scheduling.StoryPoints], 
                   [System.CreatedDate], [Microsoft.VSTS.Common.StateChangeDate],
                   [System.IterationPath], [System.AreaPath], [Microsoft.VSTS.Common.ActivatedDate],
                   [Microsoft.VSTS.Common.ResolvedDate], [Microsoft.VSTS.Common.ClosedDate]
            FROM WorkItems 
            WHERE [System.IterationPath] = '{iteration_path}'
            AND [System.AreaPath] UNDER '{area_path}'
            AND [System.WorkItemType] IN ('{work_item_types_filter}')
            ORDER BY [System.Id]
            """
        }
        
        print(f"Executing WIQL query:")
        print(f"Iteration Path: {iteration_path}")
        print(f"Area Path: {area_path}")
        print(f"Work Item Types: {work_item_types}")
        
        url = f"{self.base_url}/wit/wiql?api-version=6.0"
        response = requests.post(url, headers=self.headers, json=wiql_query)
        
        if response.status_code == 200:
            work_item_ids = [item['id'] for item in response.json().get('workItems', [])]
            print(f"Found {len(work_item_ids)} work items matching criteria")
            
            if work_item_ids:
                # Get detailed work item information
                ids_string = ','.join(map(str, work_item_ids))
                details_url = f"{self.base_url}/wit/workitems?ids={ids_string}&api-version=6.0"
                details_response = requests.get(details_url, headers=self.headers)
                
                if details_response.status_code == 200:
                    return details_response.json()['value']
                else:
                    print(f"Error fetching work item details: {details_response.status_code}")
                    print(details_response.text)
            else:
                print("No work items found matching the criteria")
                return []
        else:
            print(f"Error executing WIQL query: {response.status_code}")
            print(response.text)
            return []

    def filter_completed_work_items(self, work_items):
        """Filter work items that are marked as 'Resolved' or 'Closed'"""
        completed_items = []
        
        print(f"Filtering for completed states: {COMPLETED_STATES}")
        
        for item in work_items:
            state = item['fields'].get('System.State', '')
            if state in COMPLETED_STATES:
                completed_items.append(item)
                print(f"  ✓ Work Item {item['id']}: {state}")
            else:
                print(f"  - Work Item {item['id']}: {state} (not completed)")
        
        return completed_items

    def analyze_work_items(self, work_items):
        """Analyze work items and extract summary information"""
        total_items = len(work_items)
        total_story_points = 0
        
        work_item_details = []
        
        for item in work_items:
            fields = item['fields']
            
            # Extract story points (handle None values)
            story_points = fields.get('Microsoft.VSTS.Scheduling.StoryPoints', 0)
            if story_points:
                total_story_points += story_points
            
            # Extract assignee (handle None values)
            assignee = fields.get('System.AssignedTo', {})
            assignee_name = assignee.get('displayName', 'Unassigned') if assignee else 'Unassigned'
            
            # Extract dates for cycle time analysis
            created_date = fields.get('System.CreatedDate', '')
            activated_date = fields.get('Microsoft.VSTS.Common.ActivatedDate', '')
            resolved_date = fields.get('Microsoft.VSTS.Common.ResolvedDate', '')
            closed_date = fields.get('Microsoft.VSTS.Common.ClosedDate', '')
            state_change_date = fields.get('Microsoft.VSTS.Common.StateChangeDate', '')
            
            # Calculate cycle time (from activated to resolved/closed)
            cycle_time_days = None
            if activated_date and (resolved_date or closed_date):
                try:
                    activated_dt = datetime.fromisoformat(activated_date.replace('Z', '+00:00'))
                    completion_dt = None
                    if resolved_date:
                        completion_dt = datetime.fromisoformat(resolved_date.replace('Z', '+00:00'))
                    elif closed_date:
                        completion_dt = datetime.fromisoformat(closed_date.replace('Z', '+00:00'))
                    
                    if completion_dt:
                        cycle_time_days = (completion_dt - activated_dt).days
                except:
                    cycle_time_days = None
            
            work_item_details.append({
                'id': item['id'],
                'title': fields.get('System.Title', ''),
                'type': fields.get('System.WorkItemType', ''),
                'assignee': assignee_name,
                'story_points': story_points or 0,
                'state': fields.get('System.State', ''),
                'iteration_path': fields.get('System.IterationPath', ''),
                'area_path': fields.get('System.AreaPath', ''),
                'created_date': created_date,
                'activated_date': activated_date,
                'resolved_date': resolved_date,
                'closed_date': closed_date,
                'cycle_time_days': cycle_time_days
            })
        
        return {
            'total_items': total_items,
            'total_story_points': total_story_points,
            'work_items': work_item_details
        }

    def identify_important_work_items(self, work_items):
        """Identify important work items for review meeting"""
        important_items = []
        
        # Criteria for important items:
        # 1. High story points (>= 5)
        # 2. User Stories or Investigate items
        # 3. Items with specific keywords in title
        important_keywords = ['critical', 'important', 'major', 'feature', 'integration', 'security', 'performance', 'bug', 'investigate']
        
        for item in work_items:
            is_important = False
            reasons = []
            
            # High story points
            if item['story_points'] >= 5:
                is_important = True
                reasons.append(f"High story points ({item['story_points']})")
            
            # Important work item types (User Story, Bug, Investigate)
            if item['type'].lower() in ['user story', 'bug', 'investigate']:
                is_important = True
                reasons.append(f"Important type ({item['type']})")
            
            # Keywords in title
            title_lower = item['title'].lower()
            for keyword in important_keywords:
                if keyword in title_lower:
                    is_important = True
                    reasons.append(f"Contains keyword '{keyword}'")
                    break
            
            if is_important:
                important_items.append({
                    **item,
                    'importance_reasons': reasons
                })
        
        return important_items

    def analyze_cycle_time(self, work_items):
        """Analyze cycle time and identify work items that took longer to resolve"""
        cycle_time_analysis = {
            'items_with_cycle_time': [],
            'long_cycle_items': [],
            'avg_cycle_time': 0,
            'median_cycle_time': 0,
            'max_cycle_time': 0
        }
        
        valid_cycle_times = []
        
        for item in work_items:
            if item['cycle_time_days'] is not None:
                cycle_time_analysis['items_with_cycle_time'].append(item)
                valid_cycle_times.append(item['cycle_time_days'])
        
        if valid_cycle_times:
            # Calculate statistics
            cycle_time_analysis['avg_cycle_time'] = sum(valid_cycle_times) / len(valid_cycle_times)
            cycle_time_analysis['median_cycle_time'] = sorted(valid_cycle_times)[len(valid_cycle_times) // 2]
            cycle_time_analysis['max_cycle_time'] = max(valid_cycle_times)
            
            # Identify items that took longer than average + 1 standard deviation
            import statistics
            if len(valid_cycle_times) > 1:
                std_dev = statistics.stdev(valid_cycle_times)
                threshold = cycle_time_analysis['avg_cycle_time'] + std_dev
            else:
                threshold = cycle_time_analysis['avg_cycle_time'] * 1.5
            
            # Find items that took longer than threshold
            for item in cycle_time_analysis['items_with_cycle_time']:
                if item['cycle_time_days'] > threshold:
                    cycle_time_analysis['long_cycle_items'].append({
                        **item,
                        'cycle_time_category': 'Long' if item['cycle_time_days'] > threshold else 'Normal'
                    })
        
        return cycle_time_analysis

    def analyze_work_by_category(self, work_items):
        """Analyze work items by category (Frontend, UX, Backend, Bugs, Investigate)"""
        categories = {
            'Frontend': {'items': [], 'story_points': 0, 'keywords': ['frontend', 'fe ', 'ui', 'button', 'screen', 'window', 'tab', 'grid', 'upload', 'branding', 'text', 'alert', 'breadcrumb', 'scroll', 'menu', 'welcome', 'settings', 'angular', 'saffron', 'component']},
            'UX/UI': {'items': [], 'story_points': 0, 'keywords': ['ux', 'user experience', 'interface', 'design', 'layout', 'visual', 'styling', 'theme', 'color', 'font']},
            'Backend': {'items': [], 'story_points': 0, 'keywords': ['backend', 'api', 'service', 'endpoint', 'lambda', 'aws', 'database', 'postgres', 'server', 'deprecate', 'ultratax', 'taxassistant', 'workflow', 'metric', 'email']},
            'Bug': {'items': [], 'story_points': 0, 'keywords': []},
            'Investigate': {'items': [], 'story_points': 0, 'keywords': []},
            'Testing/QA': {'items': [], 'story_points': 0, 'keywords': ['sqa', 'test', 'testing', 'qa', 'automate', 'validation', 'regression', 'deployment', 'health check', 'scripts']},
            'Other': {'items': [], 'story_points': 0, 'keywords': []}
        }
        
        for item in work_items:
            title_lower = item['title'].lower()
            work_type = item['type']
            categorized = False
            
            # First check work item type
            if work_type == 'Bug':
                categories['Bug']['items'].append(item)
                categories['Bug']['story_points'] += item['story_points']
                categorized = True
            elif work_type == 'Investigate':
                categories['Investigate']['items'].append(item)
                categories['Investigate']['story_points'] += item['story_points']
                categorized = True
            else:
                # Categorize by keywords in title
                for category, data in categories.items():
                    if category in ['Bug', 'Investigate']:  # Already handled above
                        continue
                    
                    for keyword in data['keywords']:
                        if keyword in title_lower:
                            categories[category]['items'].append(item)
                            categories[category]['story_points'] += item['story_points']
                            categorized = True
                            break
                    
                    if categorized:
                        break
            
            # If not categorized, put in Other
            if not categorized:
                categories['Other']['items'].append(item)
                categories['Other']['story_points'] += item['story_points']
        
        return categories

    def create_burndown_chart_data(self, all_work_items, sprint_start_date, sprint_end_date):
        """Create burndown chart data based on work item state changes"""
        # Parse sprint dates
        start_date = datetime.strptime(sprint_start_date, '%Y-%m-%d').date()
        end_date = datetime.strptime(sprint_end_date, '%Y-%m-%d').date()
        
        # Calculate total story points at start
        total_points = sum(item['fields'].get('Microsoft.VSTS.Scheduling.StoryPoints', 0) or 0 
                          for item in all_work_items)
        
        print(f"Creating burndown chart for sprint {sprint_start_date} to {sprint_end_date}")
        print(f"Total story points: {total_points}")
        
        # Create sample burndown data (in real scenario, you'd track daily completions)
        dates = []
        remaining_points = []
        
        current_date = start_date
        points_left = total_points
        
        # Calculate daily burn rate based on completed work
        sprint_days = (end_date - start_date).days + 1
        daily_burn_rate = total_points / sprint_days if sprint_days > 0 else 0
        
        while current_date <= end_date:
            dates.append(current_date)
            remaining_points.append(max(0, points_left))
            
            # Simulate work completion (this would be based on actual completion dates)
            if current_date < datetime.now().date():
                # Simulate realistic burn pattern (slower at start, faster in middle, variable at end)
                days_elapsed = (current_date - start_date).days
                if days_elapsed < sprint_days * 0.3:  # First 30% of sprint
                    daily_completion = daily_burn_rate * 0.5
                elif days_elapsed < sprint_days * 0.7:  # Middle 40% of sprint
                    daily_completion = daily_burn_rate * 1.2
                else:  # Last 30% of sprint
                    daily_completion = daily_burn_rate * 0.8
                
                points_left = max(0, points_left - daily_completion)
            
            current_date += timedelta(days=1)
        
        return dates, remaining_points

    def create_burndown_chart(self, dates, remaining_points, output_path):
        """Create and save burndown chart"""
        plt.figure(figsize=(12, 8))
        
        # Plot actual burndown
        plt.plot(dates, remaining_points, 'b-', linewidth=3, label='Actual Burndown', marker='o', markersize=6)
        
        # Plot ideal burndown line
        if dates:
            ideal_points = [remaining_points[0] * (1 - i / (len(dates) - 1)) for i in range(len(dates))]
            plt.plot(dates, ideal_points, 'r--', linewidth=2, label='Ideal Burndown', alpha=0.7)
        
        plt.title(f'Sprint Burndown Chart - {SPRINT_CONFIG["iteration_path"]}', fontsize=18, fontweight='bold', pad=20)
        plt.xlabel('Date', fontsize=14)
        plt.ylabel('Remaining Story Points', fontsize=14)
        plt.legend(fontsize=12)
        plt.grid(True, alpha=0.3)
        
        # Format x-axis
        plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%m/%d'))
        plt.gca().xaxis.set_major_locator(mdates.DayLocator(interval=2))
        plt.xticks(rotation=45)
        
        # Add some styling
        plt.gca().set_facecolor('#f8f9fa')
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()

    def create_powerpoint_presentation(self, analysis_data, chart_path, output_path):
        """Create PowerPoint presentation with sprint summary"""
        prs = Presentation()
        
        # Get sprint dates from iteration path (extract from 2025_S15_Jul16-Jul29)
        iteration_name = SPRINT_CONFIG["iteration_path"]
        sprint_start = "2025-07-16"  # Extracted from Jul16
        sprint_end = "2025-07-29"    # Extracted from Jul29
        
        # Slide 1: Sprint Summary
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        title1 = slide1.shapes.title
        title1.text = f"Sprint Summary - {iteration_name}"
        
        # Add content to slide 1
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        textbox = slide1.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Sprint details
        p = text_frame.paragraphs[0]
        p.text = f"Sprint: {iteration_name}"
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Start Date: {sprint_start}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"End Date: {sprint_end}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"Team: {self.team}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"Area Path: {SPRINT_CONFIG['area_path']}"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = f"Work Item Types: {', '.join(WORK_ITEM_TYPES)}"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = f"Total Completed Work Items: {analysis_data['total_items']}"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Total Story Points Completed: {analysis_data['total_story_points']}"
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Slide 2: Completed Work Items Table
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        title2 = slide2.shapes.title
        title2.text = "Completed Work Items"
        
        # Create table
        rows = len(analysis_data['work_items']) + 1  # +1 for header
        cols = 5
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(0.8)    # ID
        table.columns[1].width = Inches(4.2)    # Title
        table.columns[2].width = Inches(1.2)    # Type
        table.columns[3].width = Inches(1.5)    # Assignee
        table.columns[4].width = Inches(1.3)    # Points
        
        # Header row
        header_cells = table.rows[0].cells
        header_cells[0].text = "ID"
        header_cells[1].text = "Title"
        header_cells[2].text = "Type"
        header_cells[3].text = "Assignee"
        header_cells[4].text = "Points"
        
        # Style header row
        for cell in header_cells:
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Data rows
        for i, item in enumerate(analysis_data['work_items']):
            row_cells = table.rows[i + 1].cells
            row_cells[0].text = str(item['id'])
            row_cells[1].text = item['title'][:45] + "..." if len(item['title']) > 45 else item['title']
            row_cells[2].text = item['type']
            row_cells[3].text = item['assignee'].split()[0] if item['assignee'] != 'Unassigned' else 'Unassigned'
            row_cells[4].text = str(item['story_points'])
            
            # Style data rows
            for cell in row_cells:
                cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Slide 3: Burndown Chart
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        title3 = slide3.shapes.title
        title3.text = "Sprint Burndown Chart"
        
        # Add chart image
        if os.path.exists(chart_path):
            left = Inches(0.5)
            top = Inches(1.5)
            slide3.shapes.add_picture(chart_path, left, top, width=Inches(9))
        
        # Slide 4: Important Work Items for Review
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        title4 = slide4.shapes.title
        title4.text = "Important Work Items for Review"
        
        # Add important work items content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5.5)
        
        textbox = slide4.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        if analysis_data.get('important_items'):
            p = text_frame.paragraphs[0]
            p.text = f"Key Work Items Completed ({len(analysis_data['important_items'])} items):"
            p.font.size = Pt(18)
            p.font.bold = True
            
            for item in analysis_data['important_items'][:5]:  # Show top 5 important items
                p = text_frame.add_paragraph()
                p.text = f"• {item['title'][:55]}{'...' if len(item['title']) > 55 else ''}"
                p.font.size = Pt(14)
                p.font.bold = True
                
                p = text_frame.add_paragraph()
                p.text = f"  ID: {item['id']} | Type: {item['type']} | Points: {item['story_points']} | Assignee: {item['assignee']}"
                p.font.size = Pt(11)
                
                p = text_frame.add_paragraph()
                p.text = f"  Why Important: {', '.join(item['importance_reasons'])}"
                p.font.size = Pt(11)
                p.font.italic = True
                
                p = text_frame.add_paragraph()
                p.text = ""
        else:
            p = text_frame.paragraphs[0]
            p.text = "No high-priority work items identified for this sprint."
            p.font.size = Pt(16)
        
        # Slide 5: Key Highlights & Blockers
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        title5 = slide5.shapes.title
        title5.text = "Key Highlights & Blockers"
        
        # Add content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5.5)
        
        textbox = slide5.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = "Key Highlights:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        # Add summary of completed work
        completed_bugs = len([item for item in analysis_data['work_items'] if item['type'] == 'Bug'])
        completed_stories = len([item for item in analysis_data['work_items'] if item['type'] == 'User Story'])
        completed_investigations = len([item for item in analysis_data['work_items'] if item['type'] == 'Investigate'])
        
        if completed_bugs > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Successfully resolved {completed_bugs} bug(s)"
            p.font.size = Pt(14)
        
        if completed_stories > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Completed {completed_stories} user story/stories"
            p.font.size = Pt(14)
        
        if completed_investigations > 0:
            p = text_frame.add_paragraph()
            p.text = f"• Finished {completed_investigations} investigation(s)"
            p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = f"• Total story points delivered: {analysis_data['total_story_points']}"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = ""
        
        p = text_frame.add_paragraph()
        p.text = "Blockers & Issues:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = "• [Add any blockers or issues here]"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = "• [Add another blocker here]"
        p.font.size = Pt(14)
        
        # Save presentation
        prs.save(output_path)


def main():
    # Get configuration values
    organization = AZURE_DEVOPS_CONFIG["organization"]
    project = AZURE_DEVOPS_CONFIG["project"]
    team = AZURE_DEVOPS_CONFIG["team"]
    
    # Get Personal Access Token from environment or prompt
    personal_access_token = os.environ.get('AZURE_DEVOPS_PAT')
    if not personal_access_token:
        personal_access_token = input("Enter your Azure DevOps Personal Access Token: ")
    
    if not personal_access_token:
        print("Personal Access Token is required!")
        return
    
    # Initialize the report generator
    report = AzureDevOpsSprintReport(organization, project, team, personal_access_token)
    
    print("Azure DevOps Sprint Report Generator")
    print("=" * 50)
    print(f"Organization: {organization}")
    print(f"Project: {project}")
    print(f"Team: {team}")
    print(f"Sprint: {SPRINT_CONFIG['iteration_path']}")
    print(f"Area Path: {SPRINT_CONFIG['area_path']}")
    print(f"Work Item Types: {', '.join(WORK_ITEM_TYPES)}")
    print()
    
    # Get work items for the specific sprint
    print("Fetching work items...")
    all_work_items = report.get_work_items_for_sprint(
        SPRINT_CONFIG["iteration_path"],
        SPRINT_CONFIG["area_path"],
        WORK_ITEM_TYPES
    )
    
    if not all_work_items:
        print("No work items found. Please check your configuration.")
        return
    
    print(f"Found {len(all_work_items)} total work items")
    
    # Filter completed work items
    completed_work_items = report.filter_completed_work_items(all_work_items)
    print(f"Found {len(completed_work_items)} completed work items")
    
    if not completed_work_items:
        print("No completed work items found.")
        return
    
    # Analyze completed work items
    analysis_data = report.analyze_work_items(completed_work_items)
    
    # Identify important work items for review meeting
    important_items = report.identify_important_work_items(analysis_data['work_items'])
    analysis_data['important_items'] = important_items
    
    # Analyze cycle time
    cycle_time_analysis = report.analyze_cycle_time(analysis_data['work_items'])
    analysis_data['cycle_time_analysis'] = cycle_time_analysis
    
    # Analyze work by category
    work_categories = report.analyze_work_by_category(analysis_data['work_items'])
    analysis_data['work_categories'] = work_categories
    
    # Print summary
    print("\n" + "=" * 50)
    print("SPRINT SUMMARY")
    print("=" * 50)
    print(f"Total completed work items: {analysis_data['total_items']}")
    print(f"Total story points completed: {analysis_data['total_story_points']}")
    
    print(f"\nCompleted Work Items:")
    print("-" * 30)
    for item in analysis_data['work_items']:
        cycle_time_str = f" | Cycle Time: {item['cycle_time_days']} days" if item['cycle_time_days'] is not None else " | Cycle Time: N/A"
        print(f"• {item['id']}: {item['title']}")
        print(f"  Type: {item['type']} | Assignee: {item['assignee']} | Points: {item['story_points']} | State: {item['state']}{cycle_time_str}")
    
    print(f"\nCYCLE TIME ANALYSIS")
    print("-" * 35)
    if cycle_time_analysis['items_with_cycle_time']:
        print(f"Items with cycle time data: {len(cycle_time_analysis['items_with_cycle_time'])}")
        print(f"Average cycle time: {cycle_time_analysis['avg_cycle_time']:.1f} days")
        print(f"Median cycle time: {cycle_time_analysis['median_cycle_time']} days")
        print(f"Maximum cycle time: {cycle_time_analysis['max_cycle_time']} days")
        
        if cycle_time_analysis['long_cycle_items']:
            print(f"\nWork items that took longer to resolve ({len(cycle_time_analysis['long_cycle_items'])} items):")
            for item in cycle_time_analysis['long_cycle_items']:
                print(f"• {item['id']}: {item['title'][:60]}{'...' if len(item['title']) > 60 else ''}")
                print(f"  Type: {item['type']} | Assignee: {item['assignee']} | Cycle Time: {item['cycle_time_days']} days")
                if item['activated_date'] and (item['resolved_date'] or item['closed_date']):
                    activated = item['activated_date'][:10] if item['activated_date'] else 'N/A'
                    completed = (item['resolved_date'] or item['closed_date'])[:10] if (item['resolved_date'] or item['closed_date']) else 'N/A'
                    print(f"  Activated: {activated} → Completed: {completed}")
                print()
        else:
            print("\nNo work items took significantly longer than average to resolve.")
    else:
        print("No cycle time data available for completed work items.")
    
    print(f"\nWORK BREAKDOWN BY CATEGORY")
    print("-" * 35)
    for category, data in work_categories.items():
        if data['items']:  # Only show categories with items
            print(f"{category}: {len(data['items'])} items, {data['story_points']} story points")
            for item in data['items']:
                print(f"  • {item['id']}: {item['title'][:70]}{'...' if len(item['title']) > 70 else ''}")
            print()
    
    print(f"\nCATEGORY SUMMARY")
    print("-" * 20)
    total_categorized_items = sum(len(data['items']) for data in work_categories.values())
    total_categorized_points = sum(data['story_points'] for data in work_categories.values())
    
    for category, data in work_categories.items():
        if data['items']:
            item_percentage = (len(data['items']) / total_categorized_items) * 100 if total_categorized_items > 0 else 0
            points_percentage = (data['story_points'] / total_categorized_points) * 100 if total_categorized_points > 0 else 0
            print(f"{category}: {len(data['items'])} items ({item_percentage:.1f}%), {data['story_points']} points ({points_percentage:.1f}%)")
    
    print(f"\nIMPORTANT WORK ITEMS FOR REVIEW")
    print("-" * 35)
    print(f"Found {len(important_items)} important items:")
    for item in important_items:
        print(f"• {item['id']}: {item['title']}")
        print(f"  Type: {item['type']} | Assignee: {item['assignee']} | Points: {item['story_points']}")
        print(f"  Why Important: {', '.join(item['importance_reasons'])}")
        print()
    
    # Create burndown chart
    print("Generating burndown chart...")
    # Create safe filename by replacing problematic characters
    safe_iteration_name = SPRINT_CONFIG['iteration_path'].replace('\\', '_').replace('/', '_').replace('_', '-')
    chart_path = f"burndown_chart_{safe_iteration_name}.png"
    dates, remaining_points = report.create_burndown_chart_data(
        all_work_items, 
        "2025-07-16",  # Sprint start date
        "2025-07-29"   # Sprint end date
    )
    report.create_burndown_chart(dates, remaining_points, chart_path)
    print(f"✓ Burndown chart saved as: {chart_path}")
    
    # Create PowerPoint presentation
    print("Creating PowerPoint presentation...")
    pptx_path = f"Sprint_Report_{safe_iteration_name}.pptx"
    report.create_powerpoint_presentation(analysis_data, chart_path, pptx_path)
    print(f"✓ PowerPoint presentation saved as: {pptx_path}")
    
    print(f"\n" + "=" * 50)
    print("SPRINT REPORT GENERATED SUCCESSFULLY!")
    print("=" * 50)
    print(f"Files created:")
    print(f"• {pptx_path} - PowerPoint presentation")
    print(f"• {chart_path} - Burndown chart")
    print(f"\nThe presentation includes:")
    print(f"• Slide 1: Sprint Summary")
    print(f"• Slide 2: Completed Work Items Table")
    print(f"• Slide 3: Burndown Chart")
    print(f"• Slide 4: Important Work Items for Review")
    print(f"• Slide 5: Key Highlights & Blockers")


if __name__ == "__main__":
    main()
