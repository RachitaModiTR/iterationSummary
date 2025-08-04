#!/usr/bin/env python3
"""
Demo Azure DevOps Sprint Report Generator
This script demonstrates the Azure DevOps sprint report functionality with sample data.
"""

import os
import sys
from datetime import datetime, timedelta
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
from pptx.util import Inches, Pt

class DemoAzureDevOpsSprintReport:
    def __init__(self):
        self.organization = "tr-tax"
        self.project = "TaxProf"
        self.team = "ADGE-Prep"
        
    def get_demo_work_items(self):
        """Generate demo work items data"""
        demo_items = [
            {
                'id': 12345,
                'title': 'Implement user authentication system',
                'type': 'Feature',
                'assignee': 'John Smith',
                'story_points': 8,
                'state': 'Resolved'
            },
            {
                'id': 12346,
                'title': 'Fix critical security vulnerability in API',
                'type': 'Bug',
                'assignee': 'Sarah Johnson',
                'story_points': 5,
                'state': 'Closed'
            },
            {
                'id': 12347,
                'title': 'Create user dashboard interface',
                'type': 'User Story',
                'assignee': 'Mike Davis',
                'story_points': 13,
                'state': 'Resolved'
            },
            {
                'id': 12348,
                'title': 'Optimize database performance queries',
                'type': 'Task',
                'assignee': 'Lisa Chen',
                'story_points': 3,
                'state': 'Closed'
            },
            {
                'id': 12349,
                'title': 'Integration with third-party payment system',
                'type': 'Feature',
                'assignee': 'David Wilson',
                'story_points': 21,
                'state': 'Resolved'
            },
            {
                'id': 12350,
                'title': 'Update documentation for API endpoints',
                'type': 'Task',
                'assignee': 'Emma Brown',
                'story_points': 2,
                'state': 'Closed'
            },
            {
                'id': 12351,
                'title': 'Implement automated testing framework',
                'type': 'User Story',
                'assignee': 'Alex Rodriguez',
                'story_points': 8,
                'state': 'Resolved'
            }
        ]
        return demo_items

    def analyze_work_items(self, work_items):
        """Analyze work items and extract summary information"""
        total_items = len(work_items)
        total_story_points = sum(item['story_points'] for item in work_items)
        
        return {
            'total_items': total_items,
            'total_story_points': total_story_points,
            'work_items': work_items
        }

    def identify_important_work_items(self, work_items):
        """Identify important work items for review meeting"""
        important_items = []
        important_keywords = ['critical', 'important', 'major', 'feature', 'integration', 'security', 'performance']
        
        for item in work_items:
            is_important = False
            reasons = []
            
            # High story points
            if item['story_points'] >= 5:
                is_important = True
                reasons.append(f"High story points ({item['story_points']})")
            
            # Important work item types
            if item['type'].lower() in ['feature', 'user story', 'epic']:
                is_important = True
                reasons.append(f"Important type ({item['type']})")
            
            # Keywords in title
            title_lower = item['title'].lower()
            for keyword in important_keywords:
                if keyword in title_lower:
                    is_important = True
                    reasons.append(f"Contains keyword '{keyword}'")
                    break
            
            if is_important:
                important_items.append({
                    **item,
                    'importance_reasons': reasons
                })
        
        return important_items

    def create_demo_burndown_chart(self, total_points, output_path):
        """Create demo burndown chart"""
        # Create demo data for 2-week sprint
        start_date = datetime.now().date() - timedelta(days=14)
        end_date = datetime.now().date()
        
        dates = []
        remaining_points = []
        ideal_points = []
        
        current_date = start_date
        points_left = total_points
        
        # Simulate work completion over sprint
        daily_completions = [0, 3, 5, 8, 2, 0, 0, 12, 7, 9, 4, 0, 0, 6, 4]  # Weekend gaps
        
        for i, daily_completion in enumerate(daily_completions):
            dates.append(current_date)
            remaining_points.append(max(0, points_left))
            ideal_points.append(total_points * (1 - i / (len(daily_completions) - 1)))
            
            points_left -= daily_completion
            current_date += timedelta(days=1)
        
        # Create the chart
        plt.figure(figsize=(12, 8))
        
        # Plot actual burndown
        plt.plot(dates, remaining_points, 'b-', linewidth=3, label='Actual Burndown', marker='o', markersize=6)
        
        # Plot ideal burndown line
        plt.plot(dates, ideal_points, 'r--', linewidth=2, label='Ideal Burndown', alpha=0.7)
        
        plt.title('Sprint Burndown Chart - ADGE-Prep Team', fontsize=18, fontweight='bold', pad=20)
        plt.xlabel('Date', fontsize=14)
        plt.ylabel('Remaining Story Points', fontsize=14)
        plt.legend(fontsize=12)
        plt.grid(True, alpha=0.3)
        
        # Format x-axis
        plt.gca().xaxis.set_major_formatter(mdates.DateFormatter('%m/%d'))
        plt.gca().xaxis.set_major_locator(mdates.DayLocator(interval=2))
        plt.xticks(rotation=45)
        
        # Add some styling
        plt.gca().set_facecolor('#f8f9fa')
        plt.tight_layout()
        plt.savefig(output_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()

    def create_powerpoint_presentation(self, analysis_data, chart_path, output_path):
        """Create PowerPoint presentation with sprint summary"""
        prs = Presentation()
        
        # Get current sprint info
        sprint_start = (datetime.now() - timedelta(days=14)).strftime('%Y-%m-%d')
        sprint_end = datetime.now().strftime('%Y-%m-%d')
        
        # Slide 1: Sprint Summary
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        title1 = slide1.shapes.title
        title1.text = "Sprint Summary - ADGE-Prep Team"
        
        # Add content to slide 1
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        textbox = slide1.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # Sprint details
        p = text_frame.paragraphs[0]
        p.text = f"Sprint: Current Sprint (Demo)"
        p.font.size = Pt(20)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Start Date: {sprint_start}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"End Date: {sprint_end}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"Team: {self.team}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"Total Completed Work Items: {analysis_data['total_items']}"
        p.font.size = Pt(16)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Total Story Points Completed: {analysis_data['total_story_points']}"
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Slide 2: Completed Work Items Table
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        title2 = slide2.shapes.title
        title2.text = "Completed Work Items"
        
        # Create table
        rows = len(analysis_data['work_items']) + 1  # +1 for header
        cols = 5
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(0.8)    # ID
        table.columns[1].width = Inches(4.2)    # Title
        table.columns[2].width = Inches(1.2)    # Type
        table.columns[3].width = Inches(1.5)    # Assignee
        table.columns[4].width = Inches(1.3)    # Points
        
        # Header row
        header_cells = table.rows[0].cells
        header_cells[0].text = "ID"
        header_cells[1].text = "Title"
        header_cells[2].text = "Type"
        header_cells[3].text = "Assignee"
        header_cells[4].text = "Points"
        
        # Style header row
        for cell in header_cells:
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Data rows
        for i, item in enumerate(analysis_data['work_items']):
            row_cells = table.rows[i + 1].cells
            row_cells[0].text = str(item['id'])
            row_cells[1].text = item['title'][:45] + "..." if len(item['title']) > 45 else item['title']
            row_cells[2].text = item['type']
            row_cells[3].text = item['assignee'].split()[0]  # First name only
            row_cells[4].text = str(item['story_points'])
            
            # Style data rows
            for cell in row_cells:
                cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Slide 3: Burndown Chart
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        title3 = slide3.shapes.title
        title3.text = "Sprint Burndown Chart"
        
        # Add chart image
        if os.path.exists(chart_path):
            left = Inches(0.5)
            top = Inches(1.5)
            slide3.shapes.add_picture(chart_path, left, top, width=Inches(9))
        
        # Slide 4: Important Work Items for Review
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        title4 = slide4.shapes.title
        title4.text = "Important Work Items for Review"
        
        # Add important work items content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5.5)
        
        textbox = slide4.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        if analysis_data.get('important_items'):
            p = text_frame.paragraphs[0]
            p.text = f"Key Work Items Completed ({len(analysis_data['important_items'])} items):"
            p.font.size = Pt(18)
            p.font.bold = True
            
            for item in analysis_data['important_items'][:5]:  # Show top 5 important items
                p = text_frame.add_paragraph()
                p.text = f"• {item['title'][:55]}{'...' if len(item['title']) > 55 else ''}"
                p.font.size = Pt(14)
                p.font.bold = True
                
                p = text_frame.add_paragraph()
                p.text = f"  ID: {item['id']} | Type: {item['type']} | Points: {item['story_points']} | Assignee: {item['assignee']}"
                p.font.size = Pt(11)
                
                p = text_frame.add_paragraph()
                p.text = f"  Why Important: {', '.join(item['importance_reasons'])}"
                p.font.size = Pt(11)
                p.font.italic = True
                
                p = text_frame.add_paragraph()
                p.text = ""
        else:
            p = text_frame.paragraphs[0]
            p.text = "No high-priority work items identified for this sprint."
            p.font.size = Pt(16)
        
        # Slide 5: Key Highlights & Blockers
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        title5 = slide5.shapes.title
        title5.text = "Key Highlights & Blockers"
        
        # Add content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5.5)
        
        textbox = slide5.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = "Key Highlights:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = "• Successfully completed major authentication system implementation"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = "• Resolved critical security vulnerability ahead of schedule"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = "• Delivered payment system integration with 21 story points"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = ""
        
        p = text_frame.add_paragraph()
        p.text = "Blockers & Issues:"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = "• [Add any blockers or issues here]"
        p.font.size = Pt(14)
        
        p = text_frame.add_paragraph()
        p.text = "• [Add another blocker here]"
        p.font.size = Pt(14)
        
        # Save presentation
        prs.save(output_path)

def main():
    print("Azure DevOps Sprint Report Generator - DEMO MODE")
    print("=" * 50)
    
    # Initialize the demo report generator
    report = DemoAzureDevOpsSprintReport()
    
    print(f"Organization: {report.organization}")
    print(f"Project: {report.project}")
    print(f"Team: {report.team}")
    print()
    
    # Get demo work items
    print("Generating demo work items...")
    work_items = report.get_demo_work_items()
    print(f"Generated {len(work_items)} demo work items")
    
    # Analyze work items
    analysis_data = report.analyze_work_items(work_items)
    
    # Identify important work items
    important_items = report.identify_important_work_items(work_items)
    analysis_data['important_items'] = important_items
    
    # Print summary
    print("\n" + "=" * 50)
    print("SPRINT SUMMARY")
    print("=" * 50)
    print(f"Total completed work items: {analysis_data['total_items']}")
    print(f"Total story points completed: {analysis_data['total_story_points']}")
    
    print(f"\nCompleted Work Items:")
    print("-" * 30)
    for item in analysis_data['work_items']:
        print(f"• {item['id']}: {item['title']}")
        print(f"  Type: {item['type']} | Assignee: {item['assignee']} | Points: {item['story_points']} | State: {item['state']}")
    
    print(f"\nIMPORTANT WORK ITEMS FOR REVIEW")
    print("-" * 35)
    print(f"Found {len(important_items)} important items:")
    for item in important_items:
        print(f"• {item['id']}: {item['title']}")
        print(f"  Type: {item['type']} | Assignee: {item['assignee']} | Points: {item['story_points']}")
        print(f"  Why Important: {', '.join(item['importance_reasons'])}")
        print()
    
    # Create burndown chart
    print("Generating burndown chart...")
    chart_path = "demo_burndown_chart.png"
    report.create_demo_burndown_chart(analysis_data['total_story_points'], chart_path)
    print(f"✓ Burndown chart saved as: {chart_path}")
    
    # Create PowerPoint presentation
    print("Creating PowerPoint presentation...")
    pptx_path = "Demo_Azure_DevOps_Sprint_Report.pptx"
    report.create_powerpoint_presentation(analysis_data, chart_path, pptx_path)
    print(f"✓ PowerPoint presentation saved as: {pptx_path}")
    
    print(f"\n" + "=" * 50)
    print("DEMO REPORT GENERATED SUCCESSFULLY!")
    print("=" * 50)
    print(f"Files created:")
    print(f"• {pptx_path} - PowerPoint presentation")
    print(f"• {chart_path} - Burndown chart")
    print(f"\nThe presentation includes:")
    print(f"• Slide 1: Sprint Summary")
    print(f"• Slide 2: Completed Work Items Table")
    print(f"• Slide 3: Burndown Chart")
    print(f"• Slide 4: Important Work Items for Review")
    print(f"• Slide 5: Key Highlights & Blockers")
    print(f"\nThis demonstrates how the system would work with real Azure DevOps data!")

if __name__ == "__main__":
    main()
